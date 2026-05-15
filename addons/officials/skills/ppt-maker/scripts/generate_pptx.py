#!/usr/bin/env python3
"""
generate_pptx.py — Generate .pptx from a JSON slide configuration.

Dependencies:
  pip install python-pptx

Usage:
  python3 generate_pptx.py --config slides.json --output output.pptx
  python3 generate_pptx.py --config slides.json --output output.pptx --template template.pptx
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Any, Optional

# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color (with or without #) to (R, G, B) tuple."""
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def apply_theme_color(hex_color: str) -> str:
    """Normalize hex color for python-pptx (no # prefix)."""
    return hex_color.lstrip("#").upper()


# ---------------------------------------------------------------------------
# PPTX generation
# ---------------------------------------------------------------------------

def create_presentation(config: dict, template_path: Optional[str] = None) -> Any:
    """Create a pptx Presentation from config dict."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        print(
            "[error] python-pptx not installed. Run: pip install python-pptx",
            file=sys.stderr,
        )
        sys.exit(1)

    if template_path and os.path.exists(template_path):
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Set 16:9 aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    theme = config.get("theme", {})
    slides_data = config.get("slides", [])

    for slide_cfg in slides_data:
        slide_type = slide_cfg.get("type", "content")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        if slide_type == "cover":
            _build_cover(slide, slide_cfg, theme)
        elif slide_type == "toc":
            _build_toc(slide, slide_cfg, theme)
        elif slide_type == "section":
            _build_section(slide, slide_cfg, theme)
        elif slide_type == "two_column":
            _build_two_column(slide, slide_cfg, theme)
        elif slide_type == "image_full":
            _build_image_full(slide, slide_cfg, theme)
        elif slide_type == "ending":
            _build_ending(slide, slide_cfg, theme)
        else:
            _build_content(slide, slide_cfg, theme)

    return prs


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _get_color(theme: dict, key: str, default: str = "000000") -> str:
    return apply_theme_color(theme.get(key, default))


def _get_font(theme: dict, key: str, default: str = "Arial") -> str:
    return theme.get(key, default)


def _add_textbox(
    slide, left, top, width, height, text: str, font_size: int, color: str,
    bold: bool = False, alignment: Any = None, font_name: str = "Arial",
    anchor: Any = None, line_spacing: float = 1.2,
):
    """Add a text box with consistent styling."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    if anchor:
        tf.paragraphs[0].alignment = alignment or PP_ALIGN.LEFT

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    try:
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
    except Exception:
        pass
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment or PP_ALIGN.LEFT
    p.line_spacing = Pt(int(font_size * line_spacing))

    return txBox


def _add_bullet_list(
    slide, left, top, width, height, items: list[str], font_size: int,
    color: str, font_name: str = "Arial",
):
    """Add a bulleted text box with multiple items."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = font_name
        try:
            p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        except Exception:
            pass
        p.level = 0
        p.line_spacing = Pt(int(font_size * 1.6))
        # Add bullet character
        p.text = f"  {item}"

    return txBox


def _add_background(slide, color: str):
    """Set slide background color."""
    from pptx.oxml.ns import qn

    bg = slide.background
    fill = bg.fill
    fill.solid()
    try:
        fill.fore_color.rgb = type(fill.fore_color.rgb)(
            *hex_to_rgb(color)
        )
    except Exception:
        pass


def _add_accent_bar(slide, left, top, width, height, color: str):
    """Add a colored rectangle as accent."""
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    try:
        shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
    except Exception:
        pass
    shape.line.fill.background()
    return shape


def _add_image_safe(slide, image_path: str, left, top, width, height):
    """Add image if file exists, otherwise add placeholder."""
    from pptx.dml.color import RGBColor

    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width, height)
            return
        except Exception:
            pass
    # Placeholder
    shape = slide.shapes.add_shape(1, left, top, width, height)  # RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(1)
    # Add placeholder text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Image]"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = 1  # center


def _build_cover(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches, Pt

    bg_color = cfg.get("background_color") or theme.get("primary_color", "1A1A2E")
    _add_background(slide, bg_color)

    primary = theme.get("primary_color", "1A1A2E")
    text_color = "FFFFFF"
    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    # Background image (optional, for AI-generated cover images)
    bg_image = cfg.get("background_image", "")
    if bg_image:
        _add_image_safe(slide, bg_image, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        # Semi-transparent overlay via dark rectangle with lower opacity would be ideal,
        # but python-pptx has limited opacity support. Use a dark shape as mask.
        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        overlay.fill.solid()
        from pptx.dml.color import RGBColor
        try:
            overlay.fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_color))
        except Exception:
            pass
        overlay.line.fill.background()
        # Set transparency via XML manipulation
        from pptx.oxml.ns import qn
        solidFill = overlay.fill._fill
        srgb = solidFill.find(qn('a:solidFill'))
        if srgb is not None:
            srgbClr = srgb[0]
            if srgbClr is not None:
                alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '40000'})
                srgbClr.append(alpha)

    # Accent line
    _add_accent_bar(
        slide, Inches(1.5), Inches(3.2), Inches(1.2), Inches(0.06),
        theme.get("accent_color", "E94560"),
    )

    # Title
    title = cfg.get("title", "")
    _add_textbox(
        slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.5),
        title, 44, text_color, bold=True, font_name=heading_font,
    )

    # Subtitle
    subtitle = cfg.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.8),
            subtitle, 20, "CCCCCC", font_name=body_font,
        )

    # Author / date
    author = cfg.get("author", "")
    if author:
        _add_textbox(
            slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
            author, 14, "999999", font_name=body_font,
        )


def _build_toc(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches

    bg_color = theme.get("background_color", "FFFFFF")
    _add_background(slide, bg_color)

    primary = _get_color(theme, "primary_color", "1A1A2E")
    text_color = _get_color(theme, "text_color", "333333")
    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    # Left accent bar
    _add_accent_bar(slide, Inches(0), Inches(0), Inches(0.1), Inches(7.5), primary)

    # Title
    title = cfg.get("title", "目录")
    _add_textbox(
        slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
        title, 32, primary, bold=True, font_name=heading_font,
    )

    # TOC items
    items = cfg.get("items", [])
    y_start = 2.0
    for i, item in enumerate(items):
        num_color = theme.get("accent_color", "E94560")
        # Number
        _add_textbox(
            slide, Inches(1.5), Inches(y_start + i * 0.9), Inches(0.8), Inches(0.6),
            f"{i + 1:02d}", 28, num_color, bold=True, font_name=heading_font,
        )
        # Text
        _add_textbox(
            slide, Inches(2.5), Inches(y_start + i * 0.9), Inches(8), Inches(0.6),
            item, 20, text_color, font_name=body_font,
        )
        # Separator line
        if i < len(items) - 1:
            _add_accent_bar(
                slide, Inches(2.5), Inches(y_start + i * 0.9 + 0.6),
                Inches(8), Inches(0.01), "E0E0E0",
            )


def _build_section(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches

    bg_color = cfg.get("background_color") or theme.get("primary_color", "1A1A2E")
    _add_background(slide, bg_color)

    primary = theme.get("primary_color", "1A1A2E")
    accent = theme.get("accent_color", "E94560")
    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    # Section number
    number = cfg.get("number", "")
    if number:
        _add_textbox(
            slide, Inches(1.5), Inches(2.0), Inches(3), Inches(1.2),
            number, 72, accent, bold=True, font_name=heading_font,
        )

    # Section title
    title = cfg.get("title", "")
    _add_textbox(
        slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.2),
        title, 40, "FFFFFF", bold=True, font_name=heading_font,
    )

    # Subtitle
    subtitle = cfg.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.8),
            subtitle, 18, "CCCCCC", font_name=body_font,
        )


def _build_content(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches

    bg_color = theme.get("background_color", "FFFFFF")
    _add_background(slide, bg_color)

    primary = _get_color(theme, "primary_color", "1A1A2E")
    text_color = _get_color(theme, "text_color", "333333")
    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    # Top accent line
    _add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), primary)

    # Title
    title = cfg.get("title", "")
    _add_textbox(
        slide, Inches(1.0), Inches(0.5), Inches(11), Inches(0.8),
        title, 28, primary, bold=True, font_name=heading_font,
    )

    # Title underline
    accent = theme.get("accent_color", primary)
    _add_accent_bar(slide, Inches(1.0), Inches(1.3), Inches(1.5), Inches(0.04), accent)

    # Image handling
    image = cfg.get("image", "")
    image_position = cfg.get("image_position", "right")
    has_image = image and os.path.exists(image)

    if has_image:
        if image_position == "right":
            text_width = Inches(6.5)
            text_left = Inches(1.0)
            img_left = Inches(8.0)
            img_top = Inches(1.8)
            img_w = Inches(4.8)
            img_h = Inches(4.8)
        elif image_position == "left":
            text_width = Inches(6.5)
            text_left = Inches(5.8)
            img_left = Inches(0.5)
            img_top = Inches(1.8)
            img_w = Inches(4.8)
            img_h = Inches(4.8)
        else:  # bottom
            text_width = Inches(11.3)
            text_left = Inches(1.0)
            img_left = Inches(1.0)
            img_top = Inches(4.5)
            img_w = Inches(11.3)
            img_h = Inches(2.5)
    else:
        text_width = Inches(11.3)
        text_left = Inches(1.0)

    # Body text
    body = cfg.get("body", "")
    y = Inches(1.8)
    if body:
        _add_textbox(
            slide, text_left, y, text_width, Inches(1.0),
            body, 16, text_color, font_name=body_font,
        )
        y += Inches(1.2)

    # Bullet points
    bullets = cfg.get("bullets", [])
    if bullets:
        _add_bullet_list(
            slide, text_left, y, text_width, Inches(4.5),
            bullets, 16, text_color, font_name=body_font,
        )

    # Image
    if has_image:
        _add_image_safe(slide, image, img_left, img_top, img_w, img_h)


def _build_two_column(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches

    bg_color = theme.get("background_color", "FFFFFF")
    _add_background(slide, bg_color)

    primary = _get_color(theme, "primary_color", "1A1A2E")
    text_color = _get_color(theme, "text_color", "333333")
    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    # Top accent line
    _add_accent_bar(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), primary)

    # Main title
    title = cfg.get("title", "")
    _add_textbox(
        slide, Inches(1.0), Inches(0.5), Inches(11), Inches(0.8),
        title, 28, primary, bold=True, font_name=heading_font,
    )

    # Divider line
    accent = theme.get("accent_color", primary)
    _add_accent_bar(slide, Inches(6.4), Inches(2.0), Inches(0.04), Inches(4.5), accent)

    # Left column
    left_title = cfg.get("left_title", "")
    _add_textbox(
        slide, Inches(1.0), Inches(2.0), Inches(5.0), Inches(0.6),
        left_title, 20, accent, bold=True, font_name=heading_font,
    )
    left_body = cfg.get("left_body", [])
    if isinstance(left_body, list):
        _add_bullet_list(
            slide, Inches(1.0), Inches(2.8), Inches(5.0), Inches(4.0),
            left_body, 15, text_color, font_name=body_font,
        )
    elif isinstance(left_body, str) and left_body:
        _add_textbox(
            slide, Inches(1.0), Inches(2.8), Inches(5.0), Inches(4.0),
            left_body, 15, text_color, font_name=body_font,
        )

    # Right column
    right_title = cfg.get("right_title", "")
    _add_textbox(
        slide, Inches(7.0), Inches(2.0), Inches(5.0), Inches(0.6),
        right_title, 20, accent, bold=True, font_name=heading_font,
    )
    right_body = cfg.get("right_body", [])
    if isinstance(right_body, list):
        _add_bullet_list(
            slide, Inches(7.0), Inches(2.8), Inches(5.0), Inches(4.0),
            right_body, 15, text_color, font_name=body_font,
        )
    elif isinstance(right_body, str) and right_body:
        _add_textbox(
            slide, Inches(7.0), Inches(2.8), Inches(5.0), Inches(4.0),
            right_body, 15, text_color, font_name=body_font,
        )


def _build_image_full(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches

    image = cfg.get("image", "")
    _add_image_safe(slide, image, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Optional overlay text at bottom
    caption = cfg.get("caption", "")
    title = cfg.get("title", "")
    bg_color = theme.get("primary_color", "1A1A2E")
    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    if title or caption:
        # Semi-transparent bar at bottom
        bar = _add_accent_bar(
            slide, Inches(0), Inches(6.0), Inches(13.333), Inches(1.5), bg_color,
        )
        from pptx.oxml.ns import qn
        solidFill = bar.fill._fill
        srgb = solidFill.find(qn('a:solidFill'))
        if srgb is not None and srgb[0] is not None:
            alpha = srgb[0].makeelement(qn('a:alpha'), {'val': '40000'})
            srgb[0].append(alpha)

        if title:
            _add_textbox(
                slide, Inches(1.0), Inches(6.2), Inches(11), Inches(0.6),
                title, 24, "FFFFFF", bold=True, font_name=heading_font,
            )
        if caption:
            _add_textbox(
                slide, Inches(1.0), Inches(6.8), Inches(11), Inches(0.5),
                caption, 14, "CCCCCC", font_name=body_font,
            )


def _build_ending(slide, cfg: dict, theme: dict) -> None:
    from pptx.util import Inches

    bg_color = cfg.get("background_color") or theme.get("primary_color", "1A1A2E")
    _add_background(slide, bg_color)

    heading_font = _get_font(theme, "heading_font", "Arial")
    body_font = _get_font(theme, "body_font", "Arial")

    title = cfg.get("title", "Thank You")
    _add_textbox(
        slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
        title, 48, "FFFFFF", bold=True, font_name=heading_font,
    )

    subtitle = cfg.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
            subtitle, 20, "CCCCCC", font_name=body_font,
        )

    # Accent line
    accent = theme.get("accent_color", "E94560")
    _add_accent_bar(
        slide, Inches(1.5), Inches(5.2), Inches(2.0), Inches(0.04), accent,
    )

    # Optional contact details
    contact = cfg.get("contact", "")
    if contact:
        _add_textbox(
            slide, Inches(1.5), Inches(5.6), Inches(10), Inches(0.6),
            contact, 14, "999999", font_name=body_font,
        )

    # QR code image
    qr_image = cfg.get("qr_image", "")
    if qr_image:
        _add_image_safe(slide, qr_image, Inches(10.5), Inches(5.0), Inches(2.0), Inches(2.0))


# ---------------------------------------------------------------------------
# Template style extraction
# ---------------------------------------------------------------------------

def extract_template_style(template_path: str) -> dict:
    """Extract theme colors and fonts from a template PPTX."""
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed"}

    try:
        prs = Presentation(template_path)
    except Exception as e:
        return {"error": str(e)}

    style: dict = {
        "colors": {},
        "fonts": {},
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slide_layouts": [],
    }

    # Extract slide layouts
    for layout in prs.slide_layouts:
        style["slide_layouts"].append({
            "name": layout.name,
            "placeholders": [
                {"idx": ph.placeholder_format.idx, "name": ph.name, "type": str(ph.placeholder_format.type)}
                for ph in layout.placeholders
            ],
        })

    return style


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    parser = argparse.ArgumentParser(description="Generate .pptx from JSON slide configuration")
    parser.add_argument("--config", default=None, help="JSON file with slide configuration")
    parser.add_argument("--output", default="output.pptx", help="Output .pptx file path")
    parser.add_argument("--template", default=None, help="Template .pptx file for theme inheritance")
    parser.add_argument("--extract-style", action="store_true", help="Extract style from template and print as JSON")
    args = parser.parse_args()

    # Style extraction mode
    if args.extract_style:
        if not args.template:
            print("[error] --extract-style requires --template", file=sys.stderr)
            sys.exit(1)
        style = extract_template_style(args.template)
        print(json.dumps(style, ensure_ascii=False, indent=2))
        if "error" in style:
            sys.exit(1)
        return

    # Load config
    if not args.config:
        print("[error] --config is required (unless using --extract-style)", file=sys.stderr)
        sys.exit(1)
    config_path = Path(args.config)
    if not config_path.exists():
        print(f"[error] Config file not found: {args.config}", file=sys.stderr)
        sys.exit(1)

    with open(config_path, "r", encoding="utf-8") as f:
        config = json.load(f)

    # Validate
    if "slides" not in config:
        print("[error] Config must contain a 'slides' array", file=sys.stderr)
        sys.exit(1)

    # Generate
    prs = create_presentation(config, args.template)

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    slide_count = len(config["slides"])
    print(f"[done] Generated {slide_count} slide(s) → {output_path}")


if __name__ == "__main__":
    main()
