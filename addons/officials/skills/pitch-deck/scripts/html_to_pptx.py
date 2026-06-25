#!/usr/bin/env python3
"""
html_to_pptx.py — Convert a pitch-deck HTML file to PPTX using slide screenshots.

Renders each <section class="slide"> in a headless browser (same rendering
pipeline as html_to_longimage.py), captures per-slide screenshots, then
assembles them into a PPTX where each slide contains the screenshot as a
full-bleed image.

This "screenshot-as-slide" approach preserves pixel-perfect visual fidelity
from the HTML rendering — CSS effects, gradients, animations, complex layouts
are all captured exactly as rendered, with zero loss.

If per-slide screenshots already exist (e.g. saved by html_to_longimage.py
with --save-slides), they can be reused via --screenshots-dir to skip
re-rendering.

Dependencies:
  pip install python-pptx Pillow playwright
  python3 -m playwright install chromium

Usage:
  # Render from HTML (auto-capture screenshots)
  python3 html_to_pptx.py <input.html> [-o output.pptx]

  # Reuse existing screenshots (skip rendering)
  python3 html_to_pptx.py <input.html> -o output.pptx --screenshots-dir ./slides/

  # Custom viewport / HiDPI
  python3 html_to_pptx.py <input.html> --width 1920 --height 1080 --scale 2.0
"""

from __future__ import annotations

import argparse
import io
import json
import os
import sys
import tempfile
from pathlib import Path
from typing import Any

# ---------------------------------------------------------------------------
# Constants (shared with html_to_longimage.py)
# ---------------------------------------------------------------------------

# Slide CSS selectors to try, in priority order
_SLIDE_SELECTORS = ("section.slide", "section[class*='slide']", ".slide", "section")

# JS to inject before screenshot: force all slides into their "visible" state
# so that IntersectionObserver-driven animations (fade-up, etc.) are fully applied.
_FORCE_VISIBLE_JS = """
// Add .visible to all slides (triggers CSS animation final states)
document.querySelectorAll('section.slide, .slide').forEach(s => s.classList.add('visible'));
// Also force any other common animation trigger classes
document.querySelectorAll('[data-animate], .animate, .reveal, .fade-in').forEach(
  el => { el.classList.add('visible'); el.classList.add('active'); el.classList.add('shown'); }
);
"""

# Delay after forcing visible state, to let CSS transitions complete
_TRANSITION_SETTLE_MS = 700


# ---------------------------------------------------------------------------
# Rendering (same pipeline as html_to_longimage.py)
# ---------------------------------------------------------------------------

def _render_slides(
    html_path: str,
    width: int = 1920,
    height: int = 1080,
    scale: float = 1.0,
) -> list[bytes]:
    """Render each slide section to PNG bytes via headless Chromium.

    Returns a list of PNG image bytes, one per slide.
    Raises RuntimeError on browser/render failures.
    """
    from playwright.sync_api import sync_playwright

    abs_html = os.path.realpath(html_path)
    file_url = f"file://{abs_html}"

    images: list[bytes] = []

    with sync_playwright() as pw:
        browser = pw.chromium.launch(headless=True)
        try:
            context = browser.new_context(
                viewport={"width": width, "height": height},
                device_scale_factor=scale,
            )
            page = context.new_page()
            page.goto(file_url, wait_until="networkidle")

            # Find the best selector that matches slides
            selector = _SLIDE_SELECTORS[0]
            slide_count = page.locator(selector).count()
            for alt in _SLIDE_SELECTORS[1:]:
                if slide_count > 0:
                    break
                selector = alt
                slide_count = page.locator(selector).count()

            if slide_count == 0:
                return []

            # Force all slides into visible/animated state before screenshot
            page.evaluate(_FORCE_VISIBLE_JS)
            page.wait_for_timeout(_TRANSITION_SETTLE_MS)

            # Capture each slide element
            for i in range(slide_count):
                el = page.locator(f"{selector} >> nth={i}")
                img_bytes = el.screenshot(type="png")
                images.append(img_bytes)
        finally:
            browser.close()

    return images


# ---------------------------------------------------------------------------
# Screenshot loading
# ---------------------------------------------------------------------------

def _load_screenshots_from_dir(screenshots_dir: str) -> list[str]:
    """Load per-slide PNG screenshots from a directory.

    Files are sorted by name (natural order: slide_01.png, slide_02.png, ...).
    Returns list of absolute file paths.
    """
    dir_path = Path(screenshots_dir)
    if not dir_path.is_dir():
        return []

    png_files = sorted(dir_path.glob("*.png"))
    return [str(f) for f in png_files]


def _save_screenshots(slide_images: list[bytes], output_dir: str, prefix: str = "slide") -> list[str]:
    """Save in-memory PNG bytes to disk and return file paths.

    Files are named: {prefix}_01.png, {prefix}_02.png, ...
    """
    os.makedirs(output_dir, exist_ok=True)
    paths: list[str] = []
    for i, img_bytes in enumerate(slide_images, start=1):
        filename = f"{prefix}_{i:02d}.png"
        filepath = os.path.join(output_dir, filename)
        with open(filepath, "wb") as f:
            f.write(img_bytes)
        paths.append(filepath)
    return paths


# ---------------------------------------------------------------------------
# PPTX assembly
# ---------------------------------------------------------------------------

def _build_pptx(
    slide_image_paths: list[str],
    output_path: str,
    slide_width_inches: float = 13.333,
    slide_height_inches: float = 7.5,
) -> str:
    """Build a PPTX file with each slide containing a full-bleed screenshot.

    Args:
        slide_image_paths: Paths to PNG screenshots, one per slide.
        output_path: Where to save the PPTX.
        slide_width_inches: Slide width in inches (default 13.333 for 16:9).
        slide_height_inches: Slide height in inches (default 7.5 for 16:9).

    Returns:
        Absolute path to the saved PPTX file.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(slide_width_inches)
    prs.slide_height = Inches(slide_height_inches)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for img_path in slide_image_paths:
        slide = prs.slides.add_slide(blank_layout)
        # Full-bleed: image covers the entire slide
        slide.shapes.add_picture(
            img_path,
            Inches(0),
            Inches(0),
            Inches(slide_width_inches),
            Inches(slide_height_inches),
        )

    prs.save(output_path)
    return os.path.abspath(output_path)


def _detect_slide_aspect(
    slide_image_paths: list[str],
    width: int,
    height: int,
) -> tuple[float, float]:
    """Detect slide dimensions in inches from the first screenshot.

    Falls back to the viewport dimensions if Pillow is unavailable.
    Returns (width_inches, height_inches).
    """
    if not slide_image_paths:
        return 13.333, 7.5

    try:
        from PIL import Image as PILImage

        with PILImage.open(slide_image_paths[0]) as im:
            img_w, img_h = im.size
        # Use the actual screenshot pixel dimensions to determine aspect ratio
        # Map to inches assuming 96 DPI (standard screen DPI)
        dpi = 96.0
        return img_w / dpi, img_h / dpi
    except Exception:
        # Fallback: derive from viewport
        dpi = 96.0
        return width / dpi, height / dpi


# ---------------------------------------------------------------------------
# Main conversion entry point
# ---------------------------------------------------------------------------

def convert(
    html_path: str,
    output_path: str | None = None,
    screenshots_dir: str | None = None,
    width: int = 1920,
    height: int = 1080,
    scale: float = 1.0,
    save_screenshots: bool = False,
) -> dict[str, Any]:
    """Convert an HTML pitch-deck to PPTX via screenshots.

    Args:
        html_path: Path to the HTML pitch-deck file.
        output_path: Output PPTX path (default: same name .pptx).
        screenshots_dir: Directory of pre-existing per-slide PNGs (skip rendering).
        width: Viewport width for rendering (default 1920).
        height: Viewport height for rendering (default 1080).
        scale: Device scale factor for HiDPI rendering (default 1.0).
        save_screenshots: If True, save per-slide PNGs alongside the PPTX.

    Returns:
        Dict with ok, pptx_file, slide_count, etc.
    """
    html_path = os.path.abspath(html_path)
    if not os.path.isfile(html_path):
        return {"ok": False, "error": f"File not found: {html_path}"}
    if not html_path.lower().endswith((".html", ".htm")):
        return {"ok": False, "error": f"Not an HTML file: {html_path}"}

    # Validate parameters
    if width < 100 or height < 100:
        return {"ok": False, "error": f"Viewport too small: {width}x{height} (min 100x100)"}
    if width > 7680 or height > 4320:
        return {"ok": False, "error": f"Viewport too large: {width}x{height} (max 7680x4320)"}
    if scale <= 0:
        return {"ok": False, "error": f"scale must be positive, got {scale}"}
    if scale > 4:
        return {"ok": False, "error": f"scale exceeds maximum (4), got {scale}"}

    html_dir = os.path.dirname(html_path)
    base_name = os.path.splitext(os.path.basename(html_path))[0]

    if not output_path:
        output_path = os.path.join(html_dir, f"{base_name}.pptx")
    else:
        output_path = os.path.abspath(output_path)

    # --- Step 1: Get per-slide screenshots ---
    reused_screenshots = False
    slide_image_paths: list[str] = []

    if screenshots_dir:
        slide_image_paths = _load_screenshots_from_dir(screenshots_dir)
        if slide_image_paths:
            reused_screenshots = True

    if not slide_image_paths:
        # Render from HTML
        try:
            slide_images = _render_slides(html_path, width=width, height=height, scale=scale)
        except Exception as e:
            return {"ok": False, "error": f"Rendering failed: {e}"}

        if not slide_images:
            return {"ok": False, "error": "No slides found in HTML (expected <section class='slide'>)"}

        # Save screenshots to disk (temp or persistent)
        if save_screenshots:
            ss_dir = os.path.join(html_dir, f"{base_name}-slides")
        else:
            ss_dir = tempfile.mkdtemp(prefix="pitchdeck_pptx_")

        slide_image_paths = _save_screenshots(slide_images, ss_dir, prefix="slide")

    # --- Step 2: Detect slide aspect ratio ---
    slide_w_in, slide_h_in = _detect_slide_aspect(slide_image_paths, width, height)

    # --- Step 3: Assemble PPTX ---
    try:
        pptx_path = _build_pptx(
            slide_image_paths,
            output_path,
            slide_width_inches=slide_w_in,
            slide_height_inches=slide_h_in,
        )
    except Exception as e:
        return {"ok": False, "error": f"PPTX assembly failed: {e}"}

    result: dict[str, Any] = {
        "ok": True,
        "html_file": html_path,
        "pptx_file": pptx_path,
        "slide_count": len(slide_image_paths),
        "slide_dimensions": f"{slide_w_in:.2f}x{slide_h_in:.2f} inches",
        "reused_screenshots": reused_screenshots,
    }

    if save_screenshots and slide_image_paths:
        result["screenshots_dir"] = os.path.dirname(slide_image_paths[0])
    elif reused_screenshots and screenshots_dir:
        result["screenshots_dir"] = os.path.abspath(screenshots_dir)

    return result


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _ensure_deps() -> None:
    """Print install hint if a dependency is missing."""
    missing = []
    try:
        import pptx  # noqa: F401
    except ImportError:
        missing.append("python-pptx")
    try:
        from PIL import Image  # noqa: F401
    except ImportError:
        missing.append("Pillow")
    try:
        from playwright.sync_api import sync_playwright  # noqa: F401
    except ImportError:
        missing.append("playwright")
    if missing:
        hint = "pip install " + " ".join(missing)
        if "playwright" in missing:
            hint += " && python3 -m playwright install chromium"
        print(
            f"Missing dependencies: {', '.join(missing)}\n"
            f"Install with: {hint}",
            file=sys.stderr,
        )
        sys.exit(1)


def main() -> None:
    _ensure_deps()

    parser = argparse.ArgumentParser(
        description="Convert a pitch-deck HTML file to PPTX using slide screenshots"
    )
    parser.add_argument("html", help="Path to the HTML pitch-deck file")
    parser.add_argument("-o", "--output", default=None, help="Output PPTX path (default: same name .pptx)")
    parser.add_argument(
        "--screenshots-dir", default=None,
        help="Directory with pre-existing per-slide PNGs (skip rendering). "
             "Files sorted by name become slide order.",
    )
    parser.add_argument("--width", type=int, default=1920, help="Viewport width in pixels (default: 1920)")
    parser.add_argument("--height", type=int, default=1080, help="Viewport height in pixels (default: 1080)")
    parser.add_argument("--scale", type=float, default=1.0, help="Device scale factor for HiDPI (default: 1.0)")
    parser.add_argument(
        "--save-screenshots", action="store_true",
        help="Save per-slide PNG screenshots alongside the PPTX",
    )
    args = parser.parse_args()

    result = convert(
        args.html,
        args.output,
        screenshots_dir=args.screenshots_dir,
        width=args.width,
        height=args.height,
        scale=args.scale,
        save_screenshots=args.save_screenshots,
    )
    print(json.dumps(result, ensure_ascii=False, indent=2))

    if not result.get("ok"):
        sys.exit(1)


if __name__ == "__main__":
    main()
