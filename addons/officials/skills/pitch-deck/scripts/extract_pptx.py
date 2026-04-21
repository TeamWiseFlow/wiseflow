#!/usr/bin/env python3
"""
extract_pptx.py — Extract text, notes, and images from a .pptx file.

Dependencies:
  pip install python-pptx Pillow

Usage:
  python3 extract_pptx.py <file.pptx> [--images-dir <dir>]

Output (JSON to stdout):
  {
    "ok": true,
    "file": "presentation.pptx",
    "slide_count": 10,
    "slides": [
      {
        "index": 1,
        "title": "Slide Title",
        "text": "Full text content of the slide",
        "notes": "Speaker notes",
        "images": ["/tmp/pptx_images/slide_01_img_0.png"]
      }
    ]
  }
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path


def extract(pptx_path: str, images_dir: str | None) -> dict:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return {
            "ok": False,
            "file": pptx_path,
            "error": "python-pptx not installed. Run: pip install python-pptx",
        }

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"ok": False, "file": pptx_path, "error": str(e)}

    if images_dir:
        os.makedirs(images_dir, exist_ok=True)

    slides_data = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        texts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text_frame.text.strip()
            if not shape_text:
                continue
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.TITLE
                title = shape_text
            else:
                if hasattr(shape, "name") and "title" in shape.name.lower() and not title:
                    title = shape_text
                else:
                    texts.append(shape_text)

        notes = ""
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text.strip()

        image_paths: list[str] = []
        if images_dir:
            img_idx = 0
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        continue
                except Exception:
                    pass
                try:
                    image = shape.image
                    ext = image.ext or "png"
                    fname = f"slide_{idx:02d}_img_{img_idx}.{ext}"
                    fpath = os.path.join(images_dir, fname)
                    with open(fpath, "wb") as f:
                        f.write(image.blob)
                    image_paths.append(fpath)
                    img_idx += 1
                except Exception:
                    continue

        slides_data.append(
            {
                "index": idx,
                "title": title,
                "text": "\n".join(texts),
                "notes": notes,
                "images": image_paths,
            }
        )

    return {
        "ok": True,
        "file": pptx_path,
        "slide_count": len(slides_data),
        "slides": slides_data,
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract content from a .pptx file")
    parser.add_argument("file", help="Path to the .pptx file")
    parser.add_argument(
        "--images-dir",
        default=None,
        help="Directory to save extracted images (skipped if not specified)",
    )
    args = parser.parse_args()

    result = extract(args.file, args.images_dir)
    print(json.dumps(result, ensure_ascii=False, indent=2))

    if not result["ok"]:
        sys.exit(1)


if __name__ == "__main__":
    main()
